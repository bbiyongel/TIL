# files 폴더에 있는 모든 PPT 파일명을 file_list 리스트에 저장하기
import os
path = "./files"
file_list = os.listdir(path)
# file_list 리스트에 있는 파일들을 하나씩 열어 텍스트를 추출하고 리스트로 저장하기
from pptx import Presentation
results = []
for file_name_raw in file_list:
    result = []
    result.append(file_name_raw)
    
    file_name = "./files/" + file_name_raw
    prs = Presentation(file_name)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                result.append(paragraph.text)
    
    results.append(result)
print(results)


# 수집한 텍스트 리스트를 엑셀에서 행 단위로 저장하기
from openpyxl import Workbook
wb = Workbook()
ws = wb.active
for i in results:
    ws.append(i)
wb.save("results.xlsx")